
from ast import Not
from pptx import Presentation
from pptx.util import Pt
from PIL import Image
import os
import func_pathInfo as info
import tkinter.messagebox as msgbox
import func_pdf as pdf
import zipfile
import shutil
 
comInfo = info.read_options()

# 성적서 번호, 입고, 완료 날짜
def Update_BasicInfo(prs, year, month, day, slideIndex, shName, company = None):
    slide=prs.slides[slideIndex]    
    month = "{0}{1}".format('0',month) if len(month) == 1 else month
    day = "{0}{1}".format('0',day) if len(day) == 1 else day

    ## 성적서 번호, 입고, 완료 날짜
    for shape in slide.shapes:
        if shape.name == shName: # ppt 표 이름으로 shape 할당
            table = shape.table

            # 성적서 번호
            cell = table.cell(0,1)
            cell.text = 'A20{0}{1}'.format(year, month)
            paragraph = cell.text_frame.paragraphs[0]
            paragraph.font.size = Pt(10.5)

            # Mask 입고 일자
            cell = table.cell(0,4)
            cell.text = '20{0}.{1}.{2}'.format(year, month,day)
            paragraph = cell.text_frame.paragraphs[0]
            paragraph.font.size = Pt(10.5)  

            # 검수 완료 일자
            cell = table.cell(1,4)
            day2 = str(int(day)+1)
            day2 = "{0}{1}".format('0',day2) if len(day2) == 1 else day2
            cell.text = '20{0}.{1}.{2}'.format(year, month,day2)
            paragraph = cell.text_frame.paragraphs[0]
            paragraph.font.size = Pt(10.5)
            
            if company != None:
                cell = table.cell(1,1)
                cell.text = company
                paragraph = cell.text_frame.paragraphs[0]
                paragraph.font.size = Pt(10.5)
                
    
    
# Mask ID 
def Update_MaskID(prs, year, month, day, slideIndex, shName, id):
    slide=prs.slides[slideIndex]    
    month = "{0}{1}".format('0',month) if len(month) == 1 else month
    day = "{0}{1}".format('0',day) if len(day) == 1 else day

    for shape in slide.shapes:
        if shape.name == shName: # ppt 표 이름으로 shape 할당
            table = shape.table

            cell = table.cell(0,1)
            cell.text = str(id)
            paragraph = cell.text_frame.paragraphs[0]
            paragraph.font.size = Pt(12)
            paragraph.font.bold = True 
    
    
def GetFilePath(directoryPath, max_files, startIndex=None, except_list = None):
    try:
        # 파일 정렬 조건 : 시간순으로, 오름차순
        dir_name = directoryPath
        max_files = max_files
        # Get list of all files only in the given directory
        list_of_files = []
        t_list_of_files = filter( lambda x: os.path.isfile(os.path.join(dir_name, x)),
                                os.listdir(dir_name) )
        
        # Sort list of files based on last modification time in ascending order
        t_list_of_files = sorted( t_list_of_files,
                                key = lambda x: os.path.getmtime(os.path.join(dir_name, x))
                                )
        
        for item in t_list_of_files:
            if item.endswith('.jpg') or item.endswith('.JPG'):
                list_of_files.append(item)
            
        start = 0 if startIndex == None else startIndex
        
        if len(list_of_files) < 1:
            msgbox.showerror("에러", f'다음 경로에 파일이 존재하지 않습니다\n확인 후 성적서를 재생성하세요.\n{dir_name}')
        
    except:
        msgbox.showerror("에러", f'경로가 존재하지 않습니다\n확인 후 성적서를 재생성하세요.\n{dir_name}')
        
    if not except_list == None:
        temp = []
        list_of_files = list_of_files[start:(start+max_files+len(except_list))]
        for i in range(len(list_of_files)):
            if not(i in except_list):
                temp.append(list_of_files[i])
        
        return temp
    else:
        return list_of_files[start:(start+max_files)]

    
def Insert_Pictures(prs, path, slideIndex, shName, columns, maxPicCnt, whiteRatio, 
                        temp = False, rot = False, startIndex=1, except_list=None):
    slide=prs.slides[slideIndex]
    files = GetFilePath(directoryPath=path, max_files=maxPicCnt, startIndex=startIndex, except_list = except_list)
    
    # 그림파일 형식 변환
    tempfiles = []
    
    if temp:
        if not os.path.exists(path+r'\temp'):
            os.makedirs(path+r'\temp')
        
        for i in range(len(files)):
            originFile = r'{0}\{1}'.format(path, files[i])
            tempFile = r'{0}\temp\temp_{1}'.format(path, files[i])
            img = Image.open(originFile)
            img.save(tempFile)
            tempfiles.append(tempFile)
            
    else:
        for i in range(len(files)):
            originFile = r'{0}\{1}'.format(path, files[i])
            tempfiles.append(originFile)
    
    # 그림 삽입
    for shape in slide.shapes:
        if shape.name == shName: # ppt 표 이름으로 shape 할당
            tableleft = shape.left
            tabletop = shape.top
            tableWidth = shape.width
            tableHeight = shape.height
            
            rows = maxPicCnt / columns
            picWidth = (tableWidth * (1-whiteRatio)) / columns            
            picHeight = (tableHeight * (1-whiteRatio)) / rows
            
            rowGap = (tableHeight - (picHeight*rows)) / (rows * 2)
            colGap = (tableWidth - (picWidth*columns)) / (columns * 2)
            
            rowGap2 = (tableHeight - (picWidth*rows)) / rows # front 사진만 적용
            
            startleft = tableleft + (colGap if columns > 1 else colGap * 2 * (-1))
            starttop = tabletop + (rowGap if columns > 1 else rowGap2 / 2)

                
            for i in range(len(tempfiles)):
                row = i % columns
                col = int(i / columns)
                left = startleft + (((colGap*2) + picWidth) * row)   
                top = starttop + (((rowGap*2) + picHeight) * col)                
                
                if rot:                    
                    pic = slide.shapes.add_picture(tempfiles[i], left=left, top=top, width=picHeight, height=picWidth)
                    pic.rotation = 90
                else:
                    pic = slide.shapes.add_picture(tempfiles[i], left=left, top=top, width=picWidth, height=picHeight)

    if temp:
        if os.path.exists(path+'\\temp'):
            remove_path = path+'\\temp'
            shutil.rmtree(remove_path)

def change_pptPicture_size(pptFrom):
    # ppt 파일 별도로 하나 저장
    # ppt 파일 확장자 zip 으로 변경
    # 압축 풀기
    # 파일명(MTZQ012X7)\ppt\media
    # 이미지 로드 후 재저장
    # 폴더 압축
    # zip 파일 확장자 ppt 로 변경

    old_name = pptFrom
    old_name_prefix = old_name.replace('.pptx','')
    new_name = old_name_prefix + '.zip'

    try:
        if os.path.isfile(new_name) :
            os.remove(new_name)
        
        os.rename(old_name,new_name)
        
    except Exception as e:
        print('예외가 발생했습니다\n',e)
        
    if os.path.exists(old_name_prefix):
        shutil.rmtree(old_name_prefix)
    
    output_unzip = zipfile.ZipFile(new_name, "r")  # "r": read 모드
    output_unzip.extractall(old_name_prefix)
    output_unzip.close()

    sub_path = 'ppt\media'
    detail_path = os.path.join(old_name_prefix, sub_path)
    for file in os.listdir(detail_path):
        if file.endswith('.jpeg') or file.endswith('.jpg'):
            file_path = os.path.join(detail_path,file)
            img = Image.open(file_path)
            # print('old size',img.size)
            # print('filename',file)
            
            ratio = 0.2 # <-원본 유지 비율
            new_img = img.resize((int(img.size[0]*ratio), (int(img.size[1]*ratio))))
            # print('new size',new_img.size)
            new_img.save(file_path)


    file_path = old_name_prefix
    owd = os.getcwd()  # 현재 working directory 를 기록해둔다
    os.chdir(file_path)  # 압축 파일 생성할 폴더로 working directory 를 이동시킨다

    zip_file = zipfile.ZipFile(old_name_prefix + ".zip", "w")
    for (path, dir, files) in os.walk(file_path):
        for file in files:
            # 상대경로를 활용하여 압축한다. (os.path.relpath)
            zip_file.write(os.path.join(os.path.relpath(path, file_path), file), compress_type=zipfile.ZIP_DEFLATED)

    zip_file.close()
    os.chdir(owd)  # 원래의 working directory 로 되돌린다

    os.rename(old_name_prefix + ".zip",old_name_prefix + ".pptx")
    shutil.rmtree(old_name_prefix)

def report_depo_A3(prs, workpath, year, month, day, id, final, except_list):
    maskid = id
    
    Update_BasicInfo(prs=prs, year=year, month=month, day=day, slideIndex=0, shName='tb_basicInfo')
    Update_MaskID(prs=prs, year=year, month=month, day=day, slideIndex=0, shName='tb_maskID', id=maskid)
    Insert_Pictures(prs=prs, path=info.get_path_front(workpath) ,slideIndex=1, shName='tb_front',
                        columns=1,maxPicCnt=1, whiteRatio=0.05, temp=True, rot = True, startIndex=1)
    Insert_Pictures(prs=prs, path=info.get_path_side(workpath),slideIndex=2, shName='tb_side',
                        columns=3,maxPicCnt=12, whiteRatio=0.2, temp=True, rot = False, startIndex=1, except_list=except_list)
    Insert_Pictures(prs=prs, path=info.get_path_inspect_depo_vision_A3(id, year, month, day),slideIndex=3, shName='tb_badsector',
                        columns=3,maxPicCnt=6, whiteRatio=0.1,temp=False, rot = False, startIndex=0)
    Insert_Pictures(prs=prs, path=info.get_path_hole_depo_vision_A3(id, year, month, day),slideIndex=3, shName='tb_alignkey_low',
                        columns=3,maxPicCnt=6, whiteRatio=0.1, temp=False, rot = False, startIndex=0)
    pathppt = info.get_path_save_depo(year, month, day, 'A3', 'PPT') + '\\' + maskid + '.pptx'
    print(pathppt)    
    prs.save(pathppt)
    
    # PDF 변환 준비 (사진 용량 줄이기)
    pathppt = info.get_path_save_depo(year, month, day, 'A3', 'PDF') + '\\' + maskid + '.pptx'
    prs.save(pathppt)
    change_pptPicture_size(pathppt)
    
    # PDF 변환
    pathpdf = info.get_path_save_depo(year, month, day, 'A3', 'PDF')
    pdf.save_pdf(pathppt, pathpdf, maskid, final)
    
def report_depo_A4(prs, workpath, year, month, day, id, final, except_list):

    maskid = id
    
    Update_BasicInfo(prs=prs, year=year, month=month, day=day, slideIndex=0, shName='tb_basicInfo')
    Update_MaskID(prs=prs, year=year, month=month, day=day, slideIndex=0, shName='tb_maskID', id=maskid)
    Insert_Pictures(prs=prs, path=info.get_path_front(workpath) ,slideIndex=1, shName='tb_front',
                        columns=1,maxPicCnt=1, whiteRatio=0.05, temp=True, rot = True, startIndex=1)
    Insert_Pictures(prs=prs, path=info.get_path_side(workpath),slideIndex=2, shName='tb_side',
                        columns=3,maxPicCnt=12, whiteRatio=0.2, temp=True, rot = False, startIndex=1, except_list=except_list)
    Insert_Pictures(prs=prs, path=info.get_path_inspect_depo_vision_A4(id, year, month, day),slideIndex=3, shName='tb_badsector',
                        columns=3,maxPicCnt=12, whiteRatio=0.1,temp=False, rot = False, startIndex=0)
    Insert_Pictures(prs=prs, path=info.get_path_hole_depo_vision_A4(id, year, month, day),slideIndex=4, shName='tb_alignkey_high',
                        columns=3,maxPicCnt=6, whiteRatio=0.1, temp=False, rot = False, startIndex=0)
    Insert_Pictures(prs=prs, path=info.get_path_hole_depo_vision_A4(id, year, month, day),slideIndex=4, shName='tb_alignkey_low',
                        columns=3,maxPicCnt=6, whiteRatio=0.1, temp=False, rot = False, startIndex=6)
    pathppt = info.get_path_save_depo(year, month, day, 'A4', 'PPT') + '\\' + maskid + '.pptx'
    print(pathppt)    
    prs.save(pathppt)
    
    # PDF 변환 준비 (사진 용량 줄이기)
    pathppt = info.get_path_save_depo(year, month, day, 'A3', 'PDF') + '\\' + maskid + '.pptx'
    prs.save(pathppt)
    change_pptPicture_size(pathppt)
    
    # PDF 변환    
    pathpdf = info.get_path_save_depo(year, month, day, 'A4', 'PDF')
    pdf.save_pdf(pathppt, pathpdf, maskid, final)
        
# def report_depo_A6(prs, workpath, year, month, day, id, final, except_list):

#     maskid = id
    
#     Update_BasicInfo(prs=prs, year=year, month=month, day=day, slideIndex=0, shName='tb_basicInfo')
#     Update_MaskID(prs=prs, year=year, month=month, day=day, slideIndex=0, shName='tb_maskID', id=maskid)
#     Insert_Pictures(prs=prs, path=info.get_path_front(workpath) ,slideIndex=1, shName='tb_front',
#                         columns=1,maxPicCnt=1, whiteRatio=0.05, temp=True, rot = True, startIndex=1)
#     Insert_Pictures(prs=prs, path=info.get_path_side(workpath),slideIndex=2, shName='tb_side',
#                         columns=3,maxPicCnt=12, whiteRatio=0.2, temp=True, rot = False, startIndex=1, except_list=except_list)
#     Insert_Pictures(prs=prs, path=info.get_path_inspect_depo_vision_A6(id, year, month, day),slideIndex=3, shName='tb_badsector',
#                         columns=3,maxPicCnt=12, whiteRatio=0.1,temp=False, rot = False, startIndex=0)
#     Insert_Pictures(prs=prs, path=info.get_path_hole_depo_vision_A6(id, year, month, day),slideIndex=4, shName='tb_alignkey_high',
#                         columns=3,maxPicCnt=6, whiteRatio=0.1, temp=False, rot = False, startIndex=0)
#     Insert_Pictures(prs=prs, path=info.get_path_hole_depo_vision_A6(id, year, month, day),slideIndex=4, shName='tb_alignkey_low',
#                         columns=3,maxPicCnt=6, whiteRatio=0.1, temp=False, rot = False, startIndex=6)
#     pathppt = info.get_path_save_depo(year, month, day, 'A6', 'PPT') + '\\' + maskid + '.pptx'
#     print(pathppt)    
#     prs.save(pathppt)
    
#     # PDF 변환 준비 (사진 용량 줄이기)
#     pathppt = info.get_path_save_depo(year, month, day, 'A3', 'PDF') + '\\' + maskid + '.pptx'
#     prs.save(pathppt)
#     change_pptPicture_size(pathppt)
    
#     # PDF 변환    
#     pathpdf = info.get_path_save_depo(year, month, day, 'A6', 'PDF')
#     pdf.save_pdf(pathppt, pathpdf, maskid, final)
    
def report_depo_A6(prs, workpath, year, month, day, id, final, except_list):
    maskid = id
    
    Update_BasicInfo(prs=prs, year=year, month=month, day=day, slideIndex=0, shName='tb_basicInfo')
    Update_MaskID(prs=prs, year=year, month=month, day=day, slideIndex=0, shName='tb_maskID', id=maskid)
    Insert_Pictures(prs=prs, path=info.get_path_front(workpath) ,slideIndex=1, shName='tb_front',
                        columns=1,maxPicCnt=1, whiteRatio=0.05, temp=True, rot = True, startIndex=0)
    Insert_Pictures(prs=prs, path=info.get_path_side(workpath),slideIndex=2, shName='tb_side',
                        columns=3,maxPicCnt=12, whiteRatio=0.2, temp=True, rot = False, startIndex=1, except_list=except_list)
    Insert_Pictures(prs=prs, path=info.get_path_inspect_depo_vision_A6(id, year, month, day),slideIndex=3, shName='tb_badsector',
                        columns=3,maxPicCnt=6, whiteRatio=0.1,temp=False, rot = False, startIndex=0)
    Insert_Pictures(prs=prs, path=info.get_path_hole_depo_vision_A6(id, year, month, day),slideIndex=3, shName='tb_alignkey_low',
                        columns=3,maxPicCnt=6, whiteRatio=0.1, temp=False, rot = False, startIndex=0)
    pathppt = info.get_path_save_depo(year, month, day, 'A6', 'PPT') + '\\' + maskid + '.pptx'
    print(pathppt)    
    prs.save(pathppt)
    
    # PDF 변환 준비 (사진 용량 줄이기)
    pathppt = info.get_path_save_depo(year, month, day, 'A6', 'PDF') + '\\' + maskid + '.pptx'
    prs.save(pathppt)
    change_pptPicture_size(pathppt)
    
    # PDF 변환
    pathpdf = info.get_path_save_depo(year, month, day, 'A6', 'PDF')
    pdf.save_pdf(pathppt, pathpdf, maskid, final)

def report_new_normal(prs, workpath, year, month, day, id, company, site, final, except_list):

    maskid = id
    
    Update_BasicInfo(prs=prs, year=year, month=month, day=day, slideIndex=0, shName='tb_basicInfo', company = company)
    Update_MaskID(prs=prs, year=year, month=month, day=day, slideIndex=0, shName='tb_maskID', id=maskid)
    Insert_Pictures(prs=prs, path=workpath ,slideIndex=1, shName='tb_front',
                        columns=1,maxPicCnt=1, whiteRatio=0.05, temp=True, rot = True, startIndex=0)
    
    pathppt = info.get_path_save_new(year, month, day, company, site, 'PPT') + '\\' + maskid + '.pptx'
    print(pathppt)    
    prs.save(pathppt)
    
    # PDF 변환 준비 (사진 용량 줄이기)
    pathppt = info.get_path_save_depo(year, month, day, 'A3', 'PDF') + '\\' + maskid + '.pptx'
    prs.save(pathppt)
    change_pptPicture_size(pathppt)
    
    # PDF 변환    
    pathpdf = info.get_path_save_new(year, month, day, company, site, 'PDF')
    pdf.save_pdf(pathppt, pathpdf, maskid, final)

def report_new_new_normal(prs, workpath, year, month, day, id, company, site, final, except_list):

    maskid = id
    
    Update_BasicInfo(prs=prs, year=year, month=month, day=day, slideIndex=0, shName='tb_basicInfo', company = company)
    Update_MaskID(prs=prs, year=year, month=month, day=day, slideIndex=0, shName='tb_maskID', id=maskid)
    Insert_Pictures(prs=prs, path=workpath ,slideIndex=1, shName='tb_front',
                        columns=1,maxPicCnt=1, whiteRatio=0.05, temp=True, rot = False, startIndex=0)
    Insert_Pictures(prs=prs, path=workpath ,slideIndex=2, shName='tb_back',
                        columns=1,maxPicCnt=1, whiteRatio=0.05, temp=True, rot = False, startIndex=1)
    
    pathppt = info.get_path_save_new(year, month, day, company, site, 'PPT') + '\\' + maskid + '.pptx'
    print(pathppt)    
    prs.save(pathppt)
    
    # PDF 변환 준비 (사진 용량 줄이기)
    pathppt = info.get_path_save_depo(year, month, day, 'A3', 'PDF') + '\\' + maskid + '.pptx'
    prs.save(pathppt)
    change_pptPicture_size(pathppt)
    
    # PDF 변환    
    pathpdf = info.get_path_save_new(year, month, day, company, site, 'PDF')
    pdf.save_pdf(pathppt, pathpdf, maskid, final)
    
def report_new_open_A3(prs, workpath, year, month, day, id, company, site, final, except_list):

    maskid = id
    
    Update_BasicInfo(prs=prs, year=year, month=month, day=day, slideIndex=0, shName='tb_basicInfo', company = company)
    Update_MaskID(prs=prs, year=year, month=month, day=day, slideIndex=0, shName='tb_maskID', id=maskid)
    Insert_Pictures(prs=prs, path=info.get_path_front(workpath) ,slideIndex=1, shName='tb_front',
                        columns=1,maxPicCnt=1, whiteRatio=0.05, temp=True, rot = True, startIndex=1)
    Insert_Pictures(prs=prs, path=info.get_path_side(workpath),slideIndex=2, shName='tb_side',
                        columns=3,maxPicCnt=12, whiteRatio=0.2, temp=True, rot = False, startIndex=1, except_list=except_list)
    Insert_Pictures(prs=prs, path=info.get_path_new_vision(year, month, day, company, id, 'inspect'),slideIndex=3, shName='tb_badsector',
                        columns=3,maxPicCnt=6, whiteRatio=0.1,temp=False, rot = False, startIndex=0, )
    Insert_Pictures(prs=prs, path=info.get_path_new_vision(year, month, day, company, id, 'hole'),slideIndex=3, shName='tb_alignkey_low',
                        columns=3,maxPicCnt=6, whiteRatio=0.1, temp=False, rot = False, startIndex=0)
    pathppt = info.get_path_save_new_openOrcvd(year, month, day, company, site, 'OPEN', 'PPT') + '\\' + maskid + '.pptx'
    print(pathppt)    
    prs.save(pathppt)
    
    # PDF 변환 준비 (사진 용량 줄이기)
    pathppt = info.get_path_save_depo(year, month, day, 'A3', 'PDF') + '\\' + maskid + '.pptx'
    prs.save(pathppt)
    change_pptPicture_size(pathppt)
    
    # PDF 변환    
    pathpdf = info.get_path_save_new_openOrcvd(year, month, day, company, site, 'OPEN', 'PDF')
    pdf.save_pdf(pathppt, pathpdf, maskid, final)
    
def report_new_open_A4(prs, workpath, year, month, day, id, company, site, final, except_list):

    maskid = id
    
    Update_BasicInfo(prs=prs, year=year, month=month, day=day, slideIndex=0, shName='tb_basicInfo', company = company)
    Update_MaskID(prs=prs, year=year, month=month, day=day, slideIndex=0, shName='tb_maskID', id=maskid)
    Insert_Pictures(prs=prs, path=info.get_path_front(workpath) ,slideIndex=1, shName='tb_front',
                        columns=1,maxPicCnt=1, whiteRatio=0.05, temp=True, rot = True, startIndex=1)
    Insert_Pictures(prs=prs, path=info.get_path_side(workpath),slideIndex=2, shName='tb_side',
                        columns=3,maxPicCnt=12, whiteRatio=0.2, temp=True, rot = False, startIndex=1, except_list=except_list)
    Insert_Pictures(prs=prs, path=info.get_path_new_vision(year, month, day, company, id, 'inspect'),slideIndex=3, shName='tb_badsector',
                        columns=3,maxPicCnt=12, whiteRatio=0.1,temp=False, rot = False, startIndex=0)
    Insert_Pictures(prs=prs, path=info.get_path_new_vision(year, month, day, company, id, 'hole'),slideIndex=4, shName='tb_alignkey_high',
                        columns=3,maxPicCnt=6, whiteRatio=0.1, temp=False, rot = False, startIndex=0)
    Insert_Pictures(prs=prs, path=info.get_path_new_vision(year, month, day, company, id, 'hole'),slideIndex=4, shName='tb_alignkey_low',
                        columns=3,maxPicCnt=6, whiteRatio=0.1, temp=False, rot = False, startIndex=6)
    pathppt = info.get_path_save_new_openOrcvd(year, month, day, company, site, 'OPEN', 'PPT') + '\\' + maskid + '.pptx'
    print(pathppt)    
    prs.save(pathppt)
    
    # PDF 변환 준비 (사진 용량 줄이기)
    pathppt = info.get_path_save_depo(year, month, day, 'A3', 'PDF') + '\\' + maskid + '.pptx'
    prs.save(pathppt)
    change_pptPicture_size(pathppt)
    
    # PDF 변환    
    pathpdf = info.get_path_save_new_openOrcvd(year, month, day, company, site, 'OPEN', 'PDF')
    pdf.save_pdf(pathppt, pathpdf, maskid, final)
    
def report_new_open_A6(prs, workpath, year, month, day, id, company, site, final, except_list):

    maskid = id
    
    Update_BasicInfo(prs=prs, year=year, month=month, day=day, slideIndex=0, shName='tb_basicInfo', company = company)
    Update_MaskID(prs=prs, year=year, month=month, day=day, slideIndex=0, shName='tb_maskID', id=maskid)
    Insert_Pictures(prs=prs, path=info.get_path_front(workpath) ,slideIndex=1, shName='tb_front',
                        columns=1,maxPicCnt=1, whiteRatio=0.05, temp=True, rot = True, startIndex=1)
    Insert_Pictures(prs=prs, path=info.get_path_side(workpath),slideIndex=2, shName='tb_side',
                        columns=3,maxPicCnt=12, whiteRatio=0.2, temp=True, rot = False, startIndex=1, except_list=except_list)
    Insert_Pictures(prs=prs, path=info.get_path_new_vision(year, month, day, company, id, 'inspect'),slideIndex=3, shName='tb_badsector',
                        columns=3,maxPicCnt=12, whiteRatio=0.1,temp=False, rot = False, startIndex=0)
    Insert_Pictures(prs=prs, path=info.get_path_new_vision(year, month, day, company, id, 'hole'),slideIndex=4, shName='tb_alignkey_high',
                        columns=3,maxPicCnt=6, whiteRatio=0.1, temp=False, rot = False, startIndex=0)
    Insert_Pictures(prs=prs, path=info.get_path_new_vision(year, month, day, company, id, 'hole'),slideIndex=4, shName='tb_alignkey_low',
                        columns=3,maxPicCnt=6, whiteRatio=0.1, temp=False, rot = False, startIndex=6)
    pathppt = info.get_path_save_new_openOrcvd(year, month, day, company, site, 'OPEN', 'PPT') + '\\' + maskid + '.pptx'
    print(pathppt)    
    prs.save(pathppt)
    
    # PDF 변환 준비 (사진 용량 줄이기)
    pathppt = info.get_path_save_depo(year, month, day, 'A3', 'PDF') + '\\' + maskid + '.pptx'
    prs.save(pathppt)
    change_pptPicture_size(pathppt)
    
    # PDF 변환    
    pathpdf = info.get_path_save_new_openOrcvd(year, month, day, company, site, 'OPEN', 'PDF')
    pdf.save_pdf(pathppt, pathpdf, maskid, final)

    
def report_new_cvd(prs, workpath, year, month, day, id, company, site, final, except_list):

    maskid = id
    
    Update_BasicInfo(prs=prs, year=year, month=month, day=day, slideIndex=0, shName='tb_basicInfo', company = company)
    Update_MaskID(prs=prs, year=year, month=month, day=day, slideIndex=0, shName='tb_maskID', id=maskid)
    Insert_Pictures(prs=prs, path=info.get_path_front(workpath) ,slideIndex=1, shName='tb_front',
                        columns=1,maxPicCnt=1, whiteRatio=0.05, temp=True, rot = True, startIndex=1)
    Insert_Pictures(prs=prs, path=info.get_path_side(workpath),slideIndex=2, shName='tb_side',
                        columns=3,maxPicCnt=12, whiteRatio=0.2, temp=True, rot = False, startIndex=1, except_list=except_list)
    Insert_Pictures(prs=prs, path=info.get_path_new_vision(year, month, day, company, id, 'inspect'),slideIndex=3, shName='tb_badsector',
                        columns=3,maxPicCnt=15, whiteRatio=0.1,temp=False, rot = False, startIndex=0)
    pathppt = info.get_path_save_new_openOrcvd(year, month, day, company, site, 'CVD', 'PPT') + '\\' + maskid + '.pptx'
    print(pathppt)    
    prs.save(pathppt)
    
    # PDF 변환 준비 (사진 용량 줄이기)
    pathppt = info.get_path_save_depo(year, month, day, 'A3', 'PDF') + '\\' + maskid + '.pptx'
    prs.save(pathppt)
    change_pptPicture_size(pathppt)
    
    # PDF 변환    
    pathpdf = info.get_path_save_new_openOrcvd(year, month, day, company, site, 'CVD', 'PDF')
    pdf.save_pdf(pathppt, pathpdf, maskid, final)
    

def extract_type(path):
    i_depo = -7
    i_depo_des = -3
        
    i_new = -8
    i_new_normal_des = -3
    i_new_normal_com = -4   
     
    i_new_open_cvd_des = -4
    i_new_open_cvd_com = -5
    i_new_open_cvd_type = -3 
    
    i_id = -1
             
    temp_list = path.split('\\')  
        
    id = temp_list[i_id]
    if temp_list[i_depo] == "1. 증착":
        if temp_list[i_depo_des] == 'A3':
            type = 'depo_A3'  
            company = None    
            site = 'A3'
            
        elif temp_list[i_depo_des] == 'A4':
            type = 'depo_A4'                 
            company = None   
            site = 'A4'

        elif temp_list[i_depo_des] == 'A6':
            type = 'depo_A6'                 
            company = None   
            site = 'A6'
        
    elif temp_list[i_new] == "2. 신규" or temp_list[i_new-1] == "2. 신규":
        # if temp_list[i_new_normal_com] == '성산' or temp_list[i_new_normal_com] == '아성':
        if temp_list[i_new_normal_com] in comInfo["NEW_NEW_NORMAL"].split(','):
            if temp_list[i_new_normal_des] == 'A6':
                type = 'new_new_normal_A6'
                company = temp_list[i_new_normal_com]    
                site = 'A6'

        elif temp_list[i_new_normal_com] in comInfo["NEW_NORMAL"].split(','):
            if temp_list[i_new_normal_des] == 'A3':
                type = 'new_normal_A3'
                company = temp_list[i_new_normal_com]  
                site = 'A3'  
                
            elif temp_list[i_new_normal_des] == 'A4':
                type = 'new_normal_A4'
                company = temp_list[i_new_normal_com]    
                site = 'A4'

            elif temp_list[i_new_normal_des] == 'A6':
                type = 'new_normal_A6'
                company = temp_list[i_new_normal_com]    
                site = 'A6'
                
                
        # elif temp_list[i_new_open_cvd_com] == '풍원' or temp_list[i_new_open_cvd_com] == '세우' or temp_list[i_new_open_cvd_com] == '핌스':
        elif temp_list[i_new_open_cvd_com] in comInfo["NEW_OPEN_CVD"].split(','):
            if temp_list[i_new_open_cvd_type] == 'OPEN':
                if temp_list[i_new_open_cvd_des] == 'A3':
                    type = 'new_open_A3'
                    company = temp_list[i_new_open_cvd_com]
                    site = 'A3'
                    
                elif temp_list[i_new_open_cvd_des] == 'A4':
                    type = 'new_open_A4'
                    company = temp_list[i_new_open_cvd_com]
                    site = 'A4'

                elif temp_list[i_new_open_cvd_des] == 'A6':
                    type = 'new_open_A6'
                    company = temp_list[i_new_open_cvd_com]
                    site = 'A6'
                    

            elif temp_list[i_new_open_cvd_type] == 'CVD':
                if temp_list[i_new_open_cvd_des] == 'A3':
                    type = 'new_cvd_A3'
                    company = temp_list[i_new_open_cvd_com]
                    site = 'A3'
                    
                elif temp_list[i_new_open_cvd_des] == 'A4':
                    type = 'new_cvd_A4'
                    company = temp_list[i_new_open_cvd_com]
                    site = 'A4'
                                
                elif temp_list[i_new_open_cvd_des] == 'A6':
                    type = 'new_cvd_A6'
                    company = temp_list[i_new_open_cvd_com]
                    site = 'A6'
            
    else:
        msgbox.showerror("에러","경로 정보가 잘못되었습니다.")
    
    return id, type, company, site

def extract_typeCount(temp_lists):
    i_depo = -7
    i_depo_des = -3
        
    i_new = -8
    i_new_normal_des = -3
    i_new_normal_com = -4   
     
    i_new_open_cvd_des = -4
    i_new_open_cvd_com = -5
    i_new_open_cvd_type = -3 
             
    
    # 0:증착 A3, 1: 증착 A4, 2: 신규(성산, 아성) A3, 3: 신규(성산, 아성) A4
    # 4:신규 오픈(핌스, 세우, 풍원) A3, 5:신규 오픈(핌스, 세우, 풍원) A4, 6:신규 CVD(핌스, 세우, 풍원) A3, 7:신규 CVD(핌스, 세우, 풍원) A4
    # 8: 증착 A6, 9: 신규(성산, 아성) A6, 10:신규 오픈(핌스, 세우, 풍원) A6, 11:신규 CVD(핌스, 세우, 풍원) A6
    
    typeCntlst = [0 for i in range(12)]
    
    for temp_list in temp_lists:
        temp_list = temp_list.split('\\')  
            
        if temp_list[i_depo] == "1. 증착":
            if temp_list[i_depo_des] == 'A3':
                typeCntlst[0] += 1
            elif temp_list[i_depo_des] == 'A4':
                typeCntlst[1] += 1
            elif temp_list[i_depo_des] == 'A6':
                typeCntlst[8] += 1

        elif temp_list[i_new] == "2. 신규" or temp_list[i_new-1] == "2. 신규":
            # if temp_list[i_new_normal_com] == '성산' or temp_list[i_new_normal_com] == '아성':
            if temp_list[i_new_normal_com] in comInfo["NEW_NORMAL"].split(','):
                if temp_list[i_new_normal_des] == 'A3':
                    typeCntlst[2] += 1 
                    
                elif temp_list[i_new_normal_des] == 'A4':
                    typeCntlst[3] += 1
                                        
                elif temp_list[i_new_normal_des] == 'A6':
                    typeCntlst[9] += 1
                    
            # elif temp_list[i_new_open_cvd_com] == '풍원' or temp_list[i_new_open_cvd_com] == '세우' or temp_list[i_new_open_cvd_com] == '핌스' or temp_list[i_new_open_cvd_com] == '오럼' or temp_list[i_new_open_cvd_com] == '위폼스':
            elif temp_list[i_new_open_cvd_com] in comInfo["NEW_OPEN_CVD"].split(','):
                if temp_list[i_new_open_cvd_type] == 'OPEN':
                    if temp_list[i_new_open_cvd_des] == 'A3':
                        typeCntlst[4] += 1             

                        
                    elif temp_list[i_new_open_cvd_des] == 'A4':
                        typeCntlst[5] += 1
                                                
                    elif temp_list[i_new_open_cvd_des] == 'A6':
                        typeCntlst[10] += 1
                        

                elif temp_list[i_new_open_cvd_type] == 'CVD':
                    if temp_list[i_new_open_cvd_des] == 'A3':
                        typeCntlst[6] += 1
                        
                    elif temp_list[i_new_open_cvd_des] == 'A4':
                        typeCntlst[7] += 1
                        
                    elif temp_list[i_new_open_cvd_des] == 'A6':
                        typeCntlst[11] += 1
    
    return typeCntlst
        

def make_reports(workpath, year, month, day, final, except_list):
    id, type, company, site = extract_type(workpath)
    pathSample = info.get_path_sample(type)
    prs = Presentation(pathSample)
    make_report(prs=prs, type=type, year=year, month=month, day=day, id=id, company=company,
                site = site, workpath=workpath, final=final, except_list=except_list)

def make_report(prs, type, year, month, day, id, company, site, workpath, final, except_list):
    if type == 'depo_A3':
        report_depo_A3(prs=prs, workpath = workpath, year=year, month=month, day=day, id=id, final=final, except_list=except_list)
    elif type == 'depo_A4':        
        report_depo_A4(prs=prs, workpath = workpath, year=year, month=month, day=day, id=id, final=final, except_list=except_list)        
    elif type == 'depo_A6':        
        report_depo_A6(prs=prs, workpath = workpath, year=year, month=month, day=day, id=id, final=final, except_list=except_list)
    elif type == 'new_normal_A3' or type == 'new_normal_A4' or type == 'new_normal_A6':         
        report_new_normal(prs=prs, workpath = workpath, year=year, month=month, day=day, id=id, company=company, site=site, final=final, except_list=except_list)
    elif type == 'new_new_normal_A6':         
        report_new_new_normal(prs=prs, workpath = workpath, year=year, month=month, day=day, id=id, company=company, site=site, final=final, except_list=except_list)
    elif type == 'new_open_A3':        
        report_new_open_A3(prs=prs, workpath = workpath, year=year, month=month, day=day, id=id, company=company, site=site, final=final, except_list=except_list)
    elif type == 'new_open_A4':
        report_new_open_A4(prs=prs, workpath = workpath, year=year, month=month, day=day, id=id, company=company, site=site, final=final, except_list=except_list)
    elif type == 'new_open_A6':
        report_new_open_A6(prs=prs, workpath = workpath, year=year, month=month, day=day, id=id, company=company, site=site, final=final, except_list=except_list)
    elif type == 'new_cvd_A3' or type == 'new_cvd_A4' or type == 'new_cvd_A6':
        report_new_cvd(prs=prs, workpath = workpath, year=year, month=month, day=day, id=id, company=company, site=site, final=final, except_list=except_list)
    else:
        msgbox.msgbox.showerror("에러", "잘못된 샘플 목록을 입력하였습니다")    
        